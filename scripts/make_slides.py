"""Generate a 5-slide deck explaining the pipeline.

Run once:  uv run python scripts/make_slides.py

Output:    slides/rx_document_processor.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ─── Brand palette ───────────────────────────────────────────────────────────
NAVY = RGBColor(0x18, 0x19, 0x3F)
TEAL = RGBColor(0x00, 0xA8, 0xB3)
TINT = RGBColor(0xD2, 0xE5, 0xE5)
SLATE = RGBColor(0x2F, 0x2E, 0x41)
INK = RGBColor(0x49, 0x50, 0x51)
SOFT_BG = RGBColor(0xF2, 0xF8, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ALERT = RGBColor(0x9C, 0x24, 0x24)


# ─── Helpers ─────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    *,
    left,
    top,
    width,
    height,
    text: str,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = NAVY,
    align=PP_ALIGN.LEFT,
    font: str = "Helvetica Neue",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_chip(
    slide, *, left, top, width, height, text: str, fill: RGBColor, font_color: RGBColor, size: int = 10
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.adjustments[0] = 0.5  # rounded
    tf = shape.text_frame
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = font_color
    return shape


def add_filled_rect(slide, *, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_brand_header(slide, title: str, subtitle: str | None = None):
    """Top brand bar with the project pill, slide title, and an optional kicker."""
    # Top accent bar
    add_filled_rect(slide, left=Inches(0), top=Inches(0), width=Inches(13.33), height=Inches(0.16), color=TEAL)

    # Brand pill
    pill_w = Inches(2.1)
    pill_h = Inches(0.32)
    pill = add_chip(
        slide,
        left=Inches(0.5),
        top=Inches(0.45),
        width=pill_w,
        height=pill_h,
        text="RX DOC PROCESSOR",
        fill=TEAL,
        font_color=WHITE,
        size=10,
    )

    # Slide title
    add_text(
        slide,
        left=Inches(0.5),
        top=Inches(0.85),
        width=Inches(12.3),
        height=Inches(0.7),
        text=title,
        size=28,
        bold=True,
        color=NAVY,
    )

    if subtitle:
        add_text(
            slide,
            left=Inches(0.5),
            top=Inches(1.55),
            width=Inches(12.3),
            height=Inches(0.4),
            text=subtitle,
            size=14,
            color=TEAL,
            bold=True,
        )


# ─── Slide builders ──────────────────────────────────────────────────────────
def slide_1_problem(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_brand_header(
        slide,
        "The intake bottleneck",
        "Long-tail B2B pharma buyers don't run EDI. Orders arrive in three messy channels.",
    )

    # Three input channel tiles
    tile_y = Inches(2.4)
    tile_h = Inches(2.0)
    tile_w = Inches(3.5)
    gap = Inches(0.4)
    start_x = Inches(0.7)

    channels = [
        ("PDF", "Faxed purchase orders", "Hospital pharmacies → 503B compounders.\nFree-text notes, table-of-line-items."),
        ("CSV", "Emailed reorder spreadsheets", "Pharmacy chains.\nColumn names chosen by the buyer:\nAcct, NDC#, item, strngth, pkg…"),
        ("SMS", "Text messages to the rep", "Independents.\n“need 20 bottles metformin\n500mg, account 12345 — Joe”"),
    ]

    for i, (chan, kicker, body) in enumerate(channels):
        x = start_x + (tile_w + gap) * i
        # Tile background
        tile = add_filled_rect(slide, left=x, top=tile_y, width=tile_w, height=tile_h, color=SOFT_BG)
        # Left teal accent bar
        add_filled_rect(slide, left=x, top=tile_y, width=Inches(0.08), height=tile_h, color=TEAL)
        # Channel chip
        add_chip(
            slide,
            left=x + Inches(0.3),
            top=tile_y + Inches(0.25),
            width=Inches(0.9),
            height=Inches(0.32),
            text=chan,
            fill=NAVY,
            font_color=WHITE,
            size=11,
        )
        # Kicker
        add_text(
            slide,
            left=x + Inches(0.3),
            top=tile_y + Inches(0.7),
            width=tile_w - Inches(0.5),
            height=Inches(0.4),
            text=kicker,
            size=14,
            bold=True,
            color=NAVY,
        )
        # Body
        add_text(
            slide,
            left=x + Inches(0.3),
            top=tile_y + Inches(1.05),
            width=tile_w - Inches(0.5),
            height=Inches(1.0),
            text=body,
            size=11,
            color=INK,
        )

    # Bottom: "today" pain bullets
    add_text(
        slide,
        left=Inches(0.7),
        top=Inches(4.7),
        width=Inches(12.0),
        height=Inches(0.4),
        text="Today: ops teams hand-key each document into the marketplace or supplier ERP.",
        size=14,
        bold=True,
        color=NAVY,
    )
    bullets = [
        "•  Slow — multi-day fulfillment lag from manual keying",
        "•  Error-prone — wrong NDC, wrong quantity → recall risk",
        "•  Engineering-bottlenecked — every new buyer channel needs a new parser",
        "•  Compliance-hostile — sensitive identifiers persist in inboxes; no consistent audit trail",
    ]
    for i, b in enumerate(bullets):
        add_text(
            slide,
            left=Inches(0.9),
            top=Inches(5.15) + Inches(0.32 * i),
            width=Inches(11.5),
            height=Inches(0.32),
            text=b,
            size=12,
            color=INK,
        )

    # Footer
    add_text(
        slide,
        left=Inches(0.5),
        top=Inches(7.0),
        width=Inches(12.3),
        height=Inches(0.3),
        text="1 / 5  ·  Order & Document Intake Pipeline",
        size=10,
        color=INK,
    )


def slide_2_pipeline(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_brand_header(
        slide,
        "Six stages, three deterministic, one LLM",
        "Bytes through one pipeline regardless of input format. New formats add no code.",
    )

    # 6 stage tiles in a vertical waterfall
    stages = [
        ("1.", "DETECT",          "Extension routing",                                 "src/detector.py"),
        ("2.", "EXTRACT",         "pypdf · pandas · passthrough",                      "src/extractors.py"),
        ("3.", "SCRUB",           "Deterministic regex → typed [REDACTED_*] tokens",   "src/scrubber.py"),
        ("4.", "STRUCTURED LLM",  "LangChain init_chat_model + with_structured_output · Claude Haiku 4.5",  "src/llm.py"),
        ("5.", "VALIDATE",        "Pydantic v2 PurchaseOrder + LineItem",              "src/schema.py"),
        ("6.", "PERSIST",         "SQLite — orders + line_items + audit_log",          "src/storage.py"),
    ]

    row_h = Inches(0.62)
    row_top = Inches(2.2)
    row_x = Inches(0.7)
    row_w = Inches(11.9)

    for i, (n, name, tech, file) in enumerate(stages):
        y = row_top + (row_h + Inches(0.08)) * i
        is_llm = i == 3

        bg = NAVY if is_llm else SOFT_BG
        accent = TEAL
        text_color = WHITE if is_llm else NAVY
        muted_color = TINT if is_llm else INK

        # Row background
        add_filled_rect(slide, left=row_x, top=y, width=row_w, height=row_h, color=bg)
        # Left accent
        add_filled_rect(slide, left=row_x, top=y, width=Inches(0.08), height=row_h, color=accent)

        # Number
        add_text(
            slide, left=row_x + Inches(0.25), top=y + Inches(0.13), width=Inches(0.5),
            height=Inches(0.4), text=n, size=18, bold=True, color=accent,
        )
        # Stage name
        add_text(
            slide, left=row_x + Inches(0.85), top=y + Inches(0.13), width=Inches(2.2),
            height=Inches(0.4), text=name, size=14, bold=True, color=text_color,
        )
        # Tech
        add_text(
            slide, left=row_x + Inches(3.1), top=y + Inches(0.15), width=Inches(7.0),
            height=Inches(0.4), text=tech, size=12, color=text_color,
        )
        # File
        add_text(
            slide, left=row_x + Inches(10.0), top=y + Inches(0.15), width=Inches(2.0),
            height=Inches(0.4), text=file, size=10, color=muted_color,
            align=PP_ALIGN.RIGHT,
        )

    # Bottom callout: LLM is the only non-deterministic stage
    add_text(
        slide,
        left=Inches(0.7),
        top=Inches(6.4),
        width=Inches(11.9),
        height=Inches(0.6),
        text=(
            "The LLM (stage 4) is the only non-deterministic stage. Compliance-relevant work — "
            "redaction, schema enforcement — stays predictable; the LLM never sees raw "
            "sensitive identifiers and cannot bypass schema validation."
        ),
        size=12,
        color=INK,
    )

    add_text(
        slide, left=Inches(0.5), top=Inches(7.0), width=Inches(12.3), height=Inches(0.3),
        text="2 / 5  ·  Architecture", size=10, color=INK,
    )


def slide_3_llm(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_brand_header(
        slide,
        "The LLM stage — universal adapter",
        "LangChain + Claude Haiku 4.5 turns any input format into one Pydantic schema.",
    )

    # Code snippet box
    add_filled_rect(
        slide, left=Inches(0.7), top=Inches(2.3), width=Inches(7.5),
        height=Inches(1.6), color=NAVY,
    )
    add_text(
        slide,
        left=Inches(0.95),
        top=Inches(2.45),
        width=Inches(7.0),
        height=Inches(0.4),
        text="src/llm.py",
        size=10,
        color=TINT,
        font="Menlo",
    )
    code = (
        "init_chat_model(\n"
        '    "claude-haiku-4-5-20251001",\n'
        '    model_provider="anthropic",\n'
        ").with_structured_output(PurchaseOrder)"
    )
    add_text(
        slide,
        left=Inches(0.95),
        top=Inches(2.8),
        width=Inches(7.0),
        height=Inches(1.1),
        text=code,
        size=13,
        color=WHITE,
        font="Menlo",
    )

    # Right column: why an LLM
    add_text(
        slide, left=Inches(8.5), top=Inches(2.3), width=Inches(4.3), height=Inches(0.4),
        text="Why an LLM here", size=15, bold=True, color=TEAL,
    )
    why = [
        "Schema-aligned extraction from heterogeneous text — PDFs, table rows, free-text SMS.",
        "New buyer channels add NO code, only sample data.",
        "Confidence + flagged_fields self-reported in the same call.",
        "Only stage that's non-deterministic; everything around it is rules-based.",
    ]
    for i, w in enumerate(why):
        add_text(
            slide,
            left=Inches(8.5),
            top=Inches(2.75) + Inches(0.45 * i),
            width=Inches(4.3),
            height=Inches(0.5),
            text=f"•  {w}",
            size=11,
            color=INK,
        )

    # Input → Output example
    add_text(
        slide, left=Inches(0.7), top=Inches(4.3), width=Inches(12.0), height=Inches(0.4),
        text="Input → output example", size=15, bold=True, color=TEAL,
    )

    # Input box
    add_filled_rect(
        slide, left=Inches(0.7), top=Inches(4.8), width=Inches(5.8),
        height=Inches(1.7), color=SOFT_BG,
    )
    add_filled_rect(
        slide, left=Inches(0.7), top=Inches(4.8), width=Inches(0.06), height=Inches(1.7),
        color=TEAL,
    )
    add_text(
        slide, left=Inches(0.85), top=Inches(4.85), width=Inches(5.5), height=Inches(0.3),
        text="SCRUBBED SMS (sent to model)", size=9, bold=True, color=TEAL,
    )
    add_text(
        slide,
        left=Inches(0.85),
        top=Inches(5.2),
        width=Inches(5.5),
        height=Inches(1.3),
        text=(
            "hi this is sarah at lakeside rx, can we get 3 boxes of\n"
            "lisinopril 20 mg, 90 ct each — patient [REDACTED_PHI]\n"
            "needs starting Monday — acct [REDACTED_ACCOUNT]"
        ),
        size=10,
        color=NAVY,
        font="Menlo",
    )

    # Output box
    add_filled_rect(
        slide, left=Inches(6.8), top=Inches(4.8), width=Inches(6.0),
        height=Inches(1.7), color=SOFT_BG,
    )
    add_filled_rect(
        slide, left=Inches(6.8), top=Inches(4.8), width=Inches(0.06), height=Inches(1.7),
        color=TEAL,
    )
    add_text(
        slide, left=Inches(6.95), top=Inches(4.85), width=Inches(5.5), height=Inches(0.3),
        text="CANONICAL PurchaseOrder (model output)", size=9, bold=True, color=TEAL,
    )
    add_text(
        slide,
        left=Inches(6.95),
        top=Inches(5.2),
        width=Inches(5.7),
        height=Inches(1.3),
        text=(
            'buyer_org_name: "Lakeside Rx"\n'
            'requested_ship_date: 2026-05-03\n'
            'line_items: [{drug: "Lisinopril", strength: "20 mg",\n'
            '              quantity: 3, unit: "box", pack: "90 ct"}]\n'
            "confidence: 0.85   flagged_fields: []"
        ),
        size=10,
        color=NAVY,
        font="Menlo",
    )

    # Cost footer
    add_text(
        slide, left=Inches(0.7), top=Inches(6.7), width=Inches(12.0), height=Inches(0.3),
        text="Cost: ~$0.001 per ingestion on Haiku 4.5. ~$0.10 to process all bundled samples.",
        size=11, color=INK,
    )

    add_text(
        slide, left=Inches(0.5), top=Inches(7.0), width=Inches(12.3), height=Inches(0.3),
        text="3 / 5  ·  The LLM stage", size=10, color=INK,
    )


def slide_4_compliance(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_brand_header(
        slide,
        "Compliance — deterministic redaction + append-only audit",
        "Sensitive identifiers never reach the model. Every ingestion leaves a record.",
    )

    # Two columns
    col_y = Inches(2.2)
    col_h = Inches(4.4)
    col_w = Inches(5.95)

    # ── Left: Scrubber ──
    add_filled_rect(slide, left=Inches(0.7), top=col_y, width=col_w, height=col_h, color=SOFT_BG)
    add_filled_rect(slide, left=Inches(0.7), top=col_y, width=Inches(0.08), height=col_h, color=TEAL)
    add_text(
        slide, left=Inches(0.95), top=col_y + Inches(0.2), width=col_w - Inches(0.4),
        height=Inches(0.4), text="SCRUBBER", size=15, bold=True, color=TEAL,
    )
    add_text(
        slide, left=Inches(0.95), top=col_y + Inches(0.6), width=col_w - Inches(0.4),
        height=Inches(0.4), text="Deterministic regex · runs BEFORE the LLM",
        size=11, bold=True, color=NAVY,
    )

    scrubs = [
        ("DEA",      r"[A-Z]{2}\d{7}",                  "[REDACTED_DEA]"),
        ("Account",  "acct/account/customer + #",       "[REDACTED_ACCOUNT]"),
        ("Routing",  "routing/aba + 9-digit",           "[REDACTED_ROUTING]"),
        ("SSN",      "###-##-####",                     "[REDACTED_SSN]"),
        ("Phone",    "US phone formats",                "[REDACTED_PHONE]"),
        ("Email",    "name@host.tld",                   "[REDACTED_EMAIL]"),
        ("PHI",      "patient/pt/for + name",           "[REDACTED_PHI]"),
    ]
    for i, (label, pat, placeholder) in enumerate(scrubs):
        y = col_y + Inches(1.15) + Inches(0.4 * i)
        add_chip(
            slide, left=Inches(0.95), top=y, width=Inches(0.85), height=Inches(0.27),
            text=label, fill=NAVY, font_color=WHITE, size=9,
        )
        add_text(
            slide, left=Inches(1.9), top=y + Inches(0.03), width=Inches(2.3),
            height=Inches(0.3), text=pat, size=10, color=INK, font="Menlo",
        )
        add_text(
            slide, left=Inches(4.3), top=y + Inches(0.03), width=Inches(2.3),
            height=Inches(0.3), text="→ " + placeholder, size=9, color=TEAL,
            font="Menlo", bold=True,
        )

    # ── Right: Audit log ──
    rcol_x = Inches(7.0)
    add_filled_rect(slide, left=rcol_x, top=col_y, width=col_w, height=col_h, color=SOFT_BG)
    add_filled_rect(slide, left=rcol_x, top=col_y, width=Inches(0.08), height=col_h, color=TEAL)
    add_text(
        slide, left=rcol_x + Inches(0.25), top=col_y + Inches(0.2), width=col_w - Inches(0.4),
        height=Inches(0.4), text="AUDIT LOG", size=15, bold=True, color=TEAL,
    )
    add_text(
        slide, left=rcol_x + Inches(0.25), top=col_y + Inches(0.6),
        width=col_w - Inches(0.4), height=Inches(0.4),
        text="Append-only · 340B-relevant procurement compliance trail",
        size=11, bold=True, color=NAVY,
    )
    audit_bullets = [
        "One row per ingestion attempt (accepted OR rejected)",
        "Captures: timestamp, source file & format, redaction signature by type",
        "LLM confidence + flagged_fields path list for review scope",
        "Status: 'accepted' / 'rejected'; FK to orders.id when accepted",
        "Persists in store.db → audit_log; survives across sessions",
        "CSV export for offline compliance review",
    ]
    for i, b in enumerate(audit_bullets):
        add_text(
            slide,
            left=rcol_x + Inches(0.25),
            top=col_y + Inches(1.15) + Inches(0.4 * i),
            width=col_w - Inches(0.4),
            height=Inches(0.4),
            text=f"•  {b}",
            size=11,
            color=INK,
        )

    # Bottom callout
    add_text(
        slide,
        left=Inches(0.7),
        top=Inches(6.75),
        width=Inches(12.0),
        height=Inches(0.4),
        text=(
            "Production swap for the scrubber: AWS Comprehend Medical or Microsoft Presidio "
            "for full HIPAA PHI detection across all 18 identifier classes."
        ),
        size=11,
        color=INK,
    )

    add_text(
        slide, left=Inches(0.5), top=Inches(7.0), width=Inches(12.3), height=Inches(0.3),
        text="4 / 5  ·  Compliance", size=10, color=INK,
    )


def slide_5_outputs(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_brand_header(
        slide,
        "Outputs and downstream wiring",
        "Three demo artifacts; six production consumers; one-line provider swap.",
    )

    # Row 1: three artifacts
    artifacts = [
        ("JSON", "Canonical PurchaseOrder", "Marketplace order-management API ingests this verbatim."),
        ("PDF", "Buyer confirmation receipt", "Brand-aligned, single-page. Email back to acknowledge."),
        ("CSV", "Audit-log export",         "Flattened audit_log for offline compliance review."),
    ]
    tile_y = Inches(2.2)
    tile_h = Inches(1.5)
    tile_w = Inches(3.95)
    start_x = Inches(0.7)
    gap = Inches(0.25)
    for i, (fmt, title, body) in enumerate(artifacts):
        x = start_x + (tile_w + gap) * i
        add_filled_rect(slide, left=x, top=tile_y, width=tile_w, height=tile_h, color=SOFT_BG)
        add_filled_rect(slide, left=x, top=tile_y, width=Inches(0.06), height=tile_h, color=TEAL)
        add_chip(
            slide, left=x + Inches(0.25), top=tile_y + Inches(0.2),
            width=Inches(0.85), height=Inches(0.32),
            text=fmt, fill=NAVY, font_color=WHITE, size=10,
        )
        add_text(
            slide, left=x + Inches(1.2), top=tile_y + Inches(0.22),
            width=tile_w - Inches(1.4), height=Inches(0.3),
            text=title, size=13, bold=True, color=NAVY,
        )
        add_text(
            slide, left=x + Inches(0.25), top=tile_y + Inches(0.65),
            width=tile_w - Inches(0.4), height=Inches(0.8),
            text=body, size=11, color=INK,
        )

    # Row 2: production wiring table
    add_text(
        slide, left=Inches(0.7), top=Inches(4.0), width=Inches(12.0), height=Inches(0.3),
        text="Production wiring (illustrative — not yet enabled)", size=14, bold=True, color=TEAL,
    )

    routes = [
        ("Marketplace OMS API",       "On accept (≥0.70 confidence)",     "REST POST · canonical JSON"),
        ("Supplier ERP",              "On accept",                         "NetSuite / SAP / Mulesoft"),
        ("Buyer confirmation email",  "On accept",                         "PDF (above), attached"),
        ("EDI 855 acknowledgement",   "If buyer is EDI-enabled",          "X12 855 over AS2 / SFTP"),
        ("Review queue UI",           "Below threshold OR flagged_fields", "Same record, reviewer scopes"),
        ("Compliance archive",        "Always",                            "Audit row → tamper-evident store"),
    ]
    row_top = Inches(4.4)
    row_h = Inches(0.32)
    row_w = Inches(11.9)
    row_x = Inches(0.7)

    for i, (where, when, fmt) in enumerate(routes):
        y = row_top + row_h * i
        add_filled_rect(
            slide, left=row_x, top=y, width=row_w, height=row_h,
            color=SOFT_BG if i % 2 == 0 else WHITE,
        )
        add_text(
            slide, left=row_x + Inches(0.2), top=y + Inches(0.05),
            width=Inches(3.6), height=Inches(0.3),
            text=where, size=11, bold=True, color=NAVY,
        )
        add_text(
            slide, left=row_x + Inches(4.0), top=y + Inches(0.05),
            width=Inches(3.8), height=Inches(0.3),
            text=when, size=11, color=INK,
        )
        add_text(
            slide, left=row_x + Inches(7.9), top=y + Inches(0.05),
            width=Inches(3.8), height=Inches(0.3),
            text=fmt, size=11, color=INK, font="Menlo",
        )

    # Bedrock swap callout
    add_filled_rect(
        slide, left=Inches(0.7), top=Inches(6.5), width=Inches(11.9), height=Inches(0.6),
        color=NAVY,
    )
    add_text(
        slide,
        left=Inches(0.95),
        top=Inches(6.6),
        width=Inches(11.5),
        height=Inches(0.5),
        text=(
            'Bedrock swap (one line):  init_chat_model("anthropic.claude-…", model_provider="bedrock_converse")'
        ),
        size=12,
        bold=True,
        color=WHITE,
        font="Menlo",
    )

    add_text(
        slide, left=Inches(0.5), top=Inches(7.2), width=Inches(12.3), height=Inches(0.3),
        text="5 / 5  ·  Outputs and production", size=10, color=INK,
    )


# ─── Entry point ─────────────────────────────────────────────────────────────
def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    slide_1_problem(prs)
    slide_2_pipeline(prs)
    slide_3_llm(prs)
    slide_4_compliance(prs)
    slide_5_outputs(prs)

    out = Path(__file__).resolve().parents[1] / "slides" / "rx_document_processor.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
